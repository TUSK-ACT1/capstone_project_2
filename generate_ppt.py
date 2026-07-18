from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_bullet_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(24)

def main():
    prs = Presentation()
    
    add_title_slide(
        prs, 
        "AI Multimodal Smart Knowledge Assistant", 
        "Capstone Project Presentation"
    )
    
    add_bullet_slide(
        prs,
        "Problem Statement",
        [
            "Develop an AI assistant supporting multiple input/output modalities.",
            "Handle text, speech, and image inputs.",
            "Generate image captions and images from text.",
            "Implement Multimodal Retrieval-Augmented Generation (RAG).",
            "Provide intelligent responses with Speech synthesis.",
            "Deliver an easy-to-use Gradio interface."
        ]
    )
    
    add_bullet_slide(
        prs,
        "Module 1: Text & Speech",
        [
            "Text Input: Users can type natural language queries.",
            "Speech Input: Speak questions, converted to text via Whisper.",
            "Text Response: LLM generates contextual answers.",
            "Speech Response: Answers converted into speech using gTTS."
        ]
    )
    
    add_bullet_slide(
        prs,
        "Module 2: Vision",
        [
            "Image Upload: Users can provide visual context.",
            "Image Captioning: Automatically describe uploaded images (BLIP/Llama Vision).",
            "Image Generation: Create novel images from text prompts (Stable Diffusion / DALL-E)."
        ]
    )
    
    add_bullet_slide(
        prs,
        "Module 3: Multimodal RAG",
        [
            "Knowledge Base: JSON-based domain knowledge (College, Tourism, Healthcare, etc.).",
            "Embeddings: Sentence transformers (MiniLM).",
            "Vector Store: FAISS for fast similarity search.",
            "Retrieval: LangChain integration for context retrieval."
        ]
    )
    
    add_bullet_slide(
        prs,
        "System Architecture & Technologies",
        [
            "User Interface: Gradio",
            "Speech-to-Text: OpenAI Whisper",
            "Text-to-Speech: gTTS",
            "Image Models: BLIP, Stable Diffusion",
            "Retrieval: LangChain, FAISS",
            "LLM: Mixtral-8x7B (via Hugging Face) or OpenAI models"
        ]
    )
    
    add_bullet_slide(
        prs,
        "Project Structure",
        [
            "app.py: Main Gradio application entry point.",
            "modules/: Dedicated modules for speech, vision, rag, and llm.",
            "knowledge_base/: JSON files for various domains.",
            "assets/: Automatically generated FAISS index and local data.",
            "Easily extensible to add new domains."
        ]
    )
    
    add_bullet_slide(
        prs,
        "Conclusion",
        [
            "A comprehensive multimodal assistant unifying Text, Speech, and Vision.",
            "Leverages state-of-the-art open source and API-based models.",
            "Flexible, easily extendable knowledge base.",
            "Ready for diverse applications (Education, Tourism, Healthcare)."
        ]
    )
    
    prs.save("Capstone_Project_Presentation.pptx")
    print("Presentation saved successfully as Capstone_Project_Presentation.pptx")

if __name__ == "__main__":
    main()
